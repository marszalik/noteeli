from pathlib import Path
import zipfile

import pytest

from app.core.config import Settings
from app.domains.workspace.service import (
    ItemAlreadyExistsError,
    InvalidPathError,
    UnsupportedFileTypeError,
    WorkspaceService,
)


def build_service(root: Path) -> WorkspaceService:
    settings = Settings(
        content_root=root,
        data_dir=root.parent / ".noteeli",
        session_secret="test-secret",
        google_client_id="",
        google_client_secret="",
    )
    return WorkspaceService(settings)


def test_build_tree_keeps_directory_hierarchy(tmp_path: Path):
    notes = tmp_path / "vault"
    nested = notes / "projekty"
    nested.mkdir(parents=True)
    (notes / "README.md").write_text("# Start\n", encoding="utf-8")
    (nested / "plan.md").write_text("## Plan\n", encoding="utf-8")
    (notes / "image.png").write_bytes(b"png")

    service = build_service(notes)
    tree = service.build_tree()

    assert tree.kind == "directory"
    assert [child.name for child in tree.children] == ["projekty", "image.png", "README.md"]
    assert tree.children[0].children[0].path == "projekty/plan.md"
    assert tree.children[1].editable is False
    assert tree.children[2].editable is True


def test_save_document_updates_markdown_file(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "todo.md"
    target.write_text("stara tresc", encoding="utf-8")

    service = build_service(notes)
    document = service.save_document("todo.md", "nowa tresc")

    assert target.read_text(encoding="utf-8") == "nowa tresc"
    assert document.content == "nowa tresc"
    assert document.editable is True


def test_path_traversal_is_blocked(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    outside = tmp_path / "outside.md"
    outside.write_text("sekret", encoding="utf-8")

    service = build_service(notes)

    with pytest.raises(InvalidPathError):
        service.read_document("../outside.md")


def test_non_markdown_file_cannot_be_saved(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    binary = notes / "archive.zip"
    binary.write_bytes(b"PK\x00\x01")

    service = build_service(notes)

    with pytest.raises(UnsupportedFileTypeError):
        service.save_document("archive.zip", "nope")


def test_unknown_text_file_opens_and_saves_as_plain_text(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "config.env"
    target.write_text("A=1\n", encoding="utf-8")

    service = build_service(notes)
    document = service.read_document("config.env")

    assert document.editable is True
    assert document.file_type == "text"
    assert document.content == "A=1\n"

    service.save_document("config.env", "A=2\n")
    assert target.read_text(encoding="utf-8") == "A=2\n"


def test_create_directory_and_markdown_file(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()

    service = build_service(notes)
    created_directory = service.create_item("", "Projekty", "directory")
    created_file = service.create_item("Projekty", "Plan", "file")

    assert created_directory.path == "Projekty"
    assert created_directory.kind == "directory"
    assert created_file.path == "Projekty/Plan.md"
    assert created_file.editable is True
    assert (notes / "Projekty").is_dir()
    assert (notes / "Projekty" / "Plan.md").is_file()


def test_create_item_rejects_duplicates(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "duplikat.md").write_text("", encoding="utf-8")

    service = build_service(notes)

    with pytest.raises(ItemAlreadyExistsError):
        service.create_item("", "duplikat.md", "file")


def test_move_item_to_another_directory(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "todo.md").write_text("", encoding="utf-8")
    (notes / "archiwum").mkdir()

    service = build_service(notes)
    moved = service.move_item("todo.md", "archiwum")

    assert moved.path == "archiwum/todo.md"
    assert moved.kind == "file"
    assert (notes / "archiwum" / "todo.md").is_file()
    assert not (notes / "todo.md").exists()


def test_move_directory_into_child_is_blocked(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "projekty" / "aktywny").mkdir(parents=True)

    service = build_service(notes)

    with pytest.raises(InvalidPathError):
        service.move_item("projekty", "projekty/aktywny")


def test_manual_order_is_persisted_in_sqlite(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "a.md").write_text("a", encoding="utf-8")
    (notes / "b.md").write_text("b", encoding="utf-8")
    (notes / "c.md").write_text("c", encoding="utf-8")

    service = build_service(notes)
    preferences = service.update_preferences(str(notes), "manual", "dark", 19)

    assert preferences.sort_mode == "manual"
    assert preferences.theme_mode == "dark"
    assert preferences.editor_font_size == 19

    service.reorder_items("", ["a.md", "c.md", "b.md"])
    tree = service.build_tree()

    assert [child.name for child in tree.children] == ["a.md", "c.md", "b.md"]


def test_browse_directories_returns_sorted_subdirectories(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "zeta").mkdir()
    (notes / "alfa").mkdir()
    (notes / "plik.md").write_text("", encoding="utf-8")

    service = build_service(notes)
    browser = service.browse_directories(str(notes))

    assert browser.current_path == str(notes.resolve())
    assert browser.parent_path == str(notes.parent.resolve())
    assert [directory.name for directory in browser.directories] == ["alfa", "zeta"]


def test_browse_directories_falls_back_to_root_for_file_path(tmp_path: Path):
    """Passing a file (not a directory) should not crash — the browser
    gracefully falls back to the configured root so the user can still
    navigate."""
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "todo.md").write_text("", encoding="utf-8")
    file_path = notes / "todo.md"

    service = build_service(notes)
    browser = service.browse_directories(str(file_path))

    # Falls back to the content root rather than raising.
    assert browser.current_path == str(notes.resolve())


def test_read_document_returns_image_preview_metadata(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    image = notes / "cover.png"
    image.write_bytes(b"png")

    service = build_service(notes)
    document = service.read_document("cover.png")

    assert document.editable is False
    assert document.previewable is True
    assert document.preview_kind == "image"


def test_read_document_returns_pdf_preview_metadata(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    pdf = notes / "guide.pdf"
    pdf.write_bytes(b"%PDF-1.7")

    service = build_service(notes)
    document = service.read_document("guide.pdf")

    assert document.editable is False
    assert document.previewable is True
    assert document.preview_kind == "pdf"


def test_resolve_embedded_asset_for_relative_image(tmp_path: Path):
    """Embedded image referenced via a sibling path inside the content
    root resolves to the underlying file."""
    notes = tmp_path / "vault"
    sub = notes / "diary"
    sub.mkdir(parents=True)
    (sub / "entry.md").write_text("![img](../assets/photo.png)", encoding="utf-8")
    asset_dir = notes / "assets"
    asset_dir.mkdir()
    (asset_dir / "photo.png").write_bytes(b"png")

    service = build_service(notes)
    resolved_path, preview_kind = service.resolve_embedded_asset(
        "diary/entry.md", "../assets/photo.png"
    )

    assert resolved_path == "assets/photo.png"
    assert preview_kind == "image"


def test_resolve_embedded_asset_uses_excalidraw_export_when_available(tmp_path: Path):
    notes = tmp_path / "vault"
    sub = notes / "diary"
    sub.mkdir(parents=True)
    source = sub / "journal.md"
    source.write_text("![[../Excalidraw/Drawing.excalidraw]]", encoding="utf-8")
    excalidraw = notes / "Excalidraw" / "Drawing.excalidraw"
    excalidraw.parent.mkdir()
    excalidraw.write_text("{}", encoding="utf-8")
    exported_image = excalidraw.with_suffix(".png")
    exported_image.write_bytes(b"png")

    service = build_service(notes)
    resolved_path, preview_kind = service.resolve_embedded_asset(
        "diary/journal.md", "../Excalidraw/Drawing.excalidraw"
    )

    assert resolved_path == "Excalidraw/Drawing.png"
    assert preview_kind == "image"


def test_upload_files_creates_multiple_items_and_skips_duplicates(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "existing.png").write_bytes(b"old")

    service = build_service(notes)
    result = service.upload_files(
        "",
        [
            ("photo.png", b"png"),
            ("existing.png", b"new"),
            ("photo.png", b"dup"),
        ],
    )

    assert [item.path for item in result.created_items] == ["photo.png"]
    assert len(result.skipped_items) == 2
    assert (notes / "photo.png").read_bytes() == b"png"
    assert (notes / "existing.png").read_bytes() == b"old"


def test_prepare_download_returns_zip_for_directory(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target_dir = notes / "assets"
    target_dir.mkdir()
    (target_dir / "cover.png").write_bytes(b"png")

    service = build_service(notes)
    archive_path, filename, is_temporary = service.prepare_download("assets")

    assert is_temporary is True
    assert filename == "assets.zip"
    assert archive_path.exists()
    with zipfile.ZipFile(archive_path) as archive:
        assert "assets/cover.png" in archive.namelist()


def test_prepare_download_returns_original_file_for_regular_file(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target_file = notes / "note.md"
    target_file.write_text("hello", encoding="utf-8")

    service = build_service(notes)
    download_path, filename, is_temporary = service.prepare_download("note.md")

    assert is_temporary is False
    assert filename == "note.md"
    assert download_path == target_file


# ── Rename ───────────────────────────────────────────────────────
# Regression cluster for the bug where the .md suffix was blindly
# appended to every renamed file, corrupting non-markdown files.

def test_rename_image_preserves_original_extension(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "before.png").write_bytes(b"png-bytes")

    service = build_service(notes)
    result = service.rename_item("before.png", "after")

    assert result.path == "after.png"
    assert (notes / "after.png").exists()
    assert not (notes / "before.png").exists()
    # bytes preserved (not silently rewrapped as a markdown file)
    assert (notes / "after.png").read_bytes() == b"png-bytes"


def test_rename_code_file_preserves_extension(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "old.py").write_text("print('hi')\n", encoding="utf-8")

    service = build_service(notes)
    result = service.rename_item("old.py", "renamed")

    assert result.path == "renamed.py"
    assert (notes / "renamed.py").read_text(encoding="utf-8") == "print('hi')\n"


def test_rename_markdown_keeps_md_suffix_when_user_omits_it(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "note.md").write_text("# hi\n", encoding="utf-8")

    service = build_service(notes)
    result = service.rename_item("note.md", "diary")

    assert result.path == "diary.md"


def test_rename_explicit_extension_replaces_original(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "draft.md").write_text("draft\n", encoding="utf-8")

    service = build_service(notes)
    result = service.rename_item("draft.md", "snippet.txt")

    assert result.path == "snippet.txt"


def test_rename_directory_does_not_get_md_suffix(tmp_path: Path):
    notes = tmp_path / "vault"
    nested = notes / "old-folder"
    nested.mkdir(parents=True)

    service = build_service(notes)
    result = service.rename_item("old-folder", "new-folder")

    assert result.path == "new-folder"
    assert (notes / "new-folder").is_dir()


def test_rename_rejects_path_separators(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "note.md").write_text("x", encoding="utf-8")

    service = build_service(notes)
    with pytest.raises(InvalidPathError):
        service.rename_item("note.md", "../escape")


def test_rename_rejects_collision_with_existing(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "a.md").write_text("a", encoding="utf-8")
    (notes / "b.md").write_text("b", encoding="utf-8")

    service = build_service(notes)
    with pytest.raises(ItemAlreadyExistsError):
        service.rename_item("a.md", "b")  # would normalize to b.md → exists


# ── Delete ───────────────────────────────────────────────────────

def test_delete_removes_file(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "scrap.md"
    target.write_text("delete me", encoding="utf-8")

    service = build_service(notes)
    service.delete_item("scrap.md")

    assert not target.exists()


def test_delete_removes_directory_recursively(tmp_path: Path):
    notes = tmp_path / "vault"
    folder = notes / "trash"
    (folder / "nested").mkdir(parents=True)
    (folder / "a.md").write_text("a", encoding="utf-8")
    (folder / "nested" / "b.md").write_text("b", encoding="utf-8")

    service = build_service(notes)
    service.delete_item("trash")

    assert not folder.exists()


def test_delete_blocks_path_traversal(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    sensitive = tmp_path / "secret.txt"
    sensitive.write_text("don't touch", encoding="utf-8")

    service = build_service(notes)
    with pytest.raises(InvalidPathError):
        service.delete_item("../secret.txt")

    # File outside the content root must be untouched.
    assert sensitive.exists()


# ── Path normalisation ──────────────────────────────────────────

def test_path_sanitisation_handles_windows_style_separators(tmp_path: Path):
    """Backslashes in the path should not crash the resolver — they get
    normalised to forward slashes so a Windows-style input still finds
    the file."""
    notes = tmp_path / "vault"
    deep = notes / "projects" / "alpha"
    deep.mkdir(parents=True)
    (deep / "spec.md").write_text("# spec\n", encoding="utf-8")

    service = build_service(notes)
    document = service.read_document("projects\\alpha/spec.md")

    # Whatever the response shape, the document content was found and read.
    assert document.content == "# spec\n"
    assert document.editable is True


def test_create_item_rejects_separator_in_name(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()

    service = build_service(notes)
    with pytest.raises(InvalidPathError):
        service.create_item("", "../evil", "file")


# ── Editor file-type classifiers ────────────────────────────────

def test_editor_file_type_classifies_markdown_json_and_text(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()

    service = build_service(notes)

    assert service.get_editor_file_type("note.md") == "markdown"
    assert service.get_editor_file_type("data.json") == "json"
    assert service.get_editor_file_type("script.py") == "text"
    assert service.is_json("payload.json") is True
    assert service.is_json("note.md") is False


def test_is_editable_true_for_markdown_and_json(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()

    service = build_service(notes)

    assert service.is_editable("a.md") is True
    assert service.is_editable("a.markdown") is True
    assert service.is_editable("a.json") is True
    # Binary preview-only files are not "editable" in this sense.
    assert service.is_editable("a.png") is False
    assert service.is_editable("a.pdf") is False


# ── Office (docx / xlsx) preview rendering ─────────────────────

def test_get_preview_kind_classifies_office_documents(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    service = build_service(notes)

    assert service.get_preview_kind("note.docx") == "docx"
    assert service.get_preview_kind("data.xlsx") == "xlsx"
    assert service.get_preview_kind("data.xlsm") == "xlsx"
    assert service.get_preview_kind("note.txt") is None


def test_render_docx_preview_returns_html(tmp_path: Path):
    """Render a tiny .docx file (built on the fly) into HTML and check
    the output contains the document's text."""
    import io
    import zipfile

    docx_xml_document = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>Hello from a docx</w:t></w:r></w:p>
    <w:p><w:r><w:t>Second paragraph</w:t></w:r></w:p>
  </w:body>
</w:document>"""
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""

    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "letter.docx"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", docx_xml_document)
    target.write_bytes(buf.getvalue())

    service = build_service(notes)
    html = service.render_office_preview("letter.docx")

    assert "<html" in html
    assert "Hello from a docx" in html
    assert "Second paragraph" in html


def test_render_xlsx_preview_returns_html_table(tmp_path: Path):
    """Render a workbook with two sheets and check both tables show up
    in the rendered HTML."""
    from openpyxl import Workbook

    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "budget.xlsx"

    wb = Workbook()
    sheet1 = wb.active
    sheet1.title = "January"
    sheet1.append(["Item", "Cost"])
    sheet1.append(["Coffee", 12.50])
    sheet1.append(["Lunch", 25])

    sheet2 = wb.create_sheet("February")
    sheet2.append(["Item", "Cost"])
    sheet2.append(["Books", 99.00])
    wb.save(str(target))

    service = build_service(notes)
    html = service.render_office_preview("budget.xlsx")

    assert "<table" in html
    assert "January" in html
    assert "February" in html
    assert "Coffee" in html
    assert "12.5" in html or "12.50" in html  # openpyxl may strip trailing zero
    assert "Books" in html


def test_render_pptx_preview_returns_html_with_slides(tmp_path: Path):
    """Build a tiny .pptx with two slides (title + bullets) and confirm
    the rendered HTML wraps each slide in its own card."""
    from pptx import Presentation
    from pptx.util import Inches

    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "deck.pptx"

    prs = Presentation()
    # Slide 1: title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Mój pierwszy slajd"
    slide1.placeholders[1].text = "Podtytuł"
    # Slide 2: title + content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Plan dnia"
    body = slide2.placeholders[1].text_frame
    body.text = "Wstać"
    body.add_paragraph().text = "Zjeść"
    body.add_paragraph().text = "Pisać kod"
    prs.save(str(target))

    service = build_service(notes)
    html = service.render_office_preview("deck.pptx")

    assert "Slajd 1" in html
    assert "Slajd 2" in html
    assert "Mój pierwszy slajd" in html
    assert "Plan dnia" in html
    assert "Wstać" in html
    assert "Zjeść" in html
    # Each slide is wrapped in its own card
    assert html.count("class='slide'") == 2 or html.count('class="slide"') == 2


def test_get_preview_kind_classifies_pptx(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    service = build_service(notes)
    assert service.get_preview_kind("deck.pptx") == "pptx"


def test_render_office_preview_rejects_non_office_file(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    (notes / "note.md").write_text("# hi\n", encoding="utf-8")

    service = build_service(notes)
    with pytest.raises(UnsupportedFileTypeError):
        service.render_office_preview("note.md")


# ── JSON document round-trip ────────────────────────────────────

def test_save_and_read_json_document(tmp_path: Path):
    notes = tmp_path / "vault"
    notes.mkdir()
    target = notes / "config.json"
    target.write_text("{}", encoding="utf-8")

    service = build_service(notes)
    saved = service.save_document("config.json", '{"theme": "dark"}')

    assert saved.file_type == "json"
    assert saved.content == '{"theme": "dark"}'
    # And re-reading round-trips.
    fresh = service.read_document("config.json")
    assert fresh.content == '{"theme": "dark"}'
